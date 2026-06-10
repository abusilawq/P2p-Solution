#!/usr/bin/env python3
"""Build the Clinic Appointment Management System presentation (16:9 PPTX).

All data, SQL and screenshots come from the team's DOCX report so the deck
matches the submitted assignment exactly.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

A = "/tmp/dms/assets"
OUT = "/tmp/dms/Clinic-Appointment-DBMS-Presentation.pptx"

# ---------- palette ----------
NAVY        = RGBColor(0x0A, 0x25, 0x40)   # deep navy (dark slides / code bg)
NAVY2       = RGBColor(0x11, 0x36, 0x5C)   # panel on dark
INK         = RGBColor(0x0F, 0x2B, 0x46)   # heading text on light
TEAL        = RGBColor(0x0D, 0x94, 0x88)   # primary accent
TEAL_BRIGHT = RGBColor(0x2D, 0xD4, 0xBF)   # accent on dark
TEAL_LIGHT  = RGBColor(0xD7, 0xF1, 0xEE)   # chip bg
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
CORAL       = RGBColor(0xE8, 0x5D, 0x5D)
BG          = RGBColor(0xF6, 0xF9, 0xFC)   # light slide bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE       = RGBColor(0x47, 0x55, 0x69)   # body text
SLATE_LIGHT = RGBColor(0x8E, 0xA3, 0xB8)   # captions
MIST        = RGBColor(0xB8, 0xC7, 0xDA)   # body text on dark
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
CHIP_GRAY   = RGBColor(0xEE, 0xF3, 0xF8)
CODE_ID     = RGBColor(0xE2, 0xE8, 0xF0)
CODE_STR    = RGBColor(0xFC, 0xD3, 0x4D)

# entity accent colours (echo the Lucidchart ERD)
C_PATIENT   = RGBColor(0x8E, 0x5B, 0xD0)
C_DOCTOR    = RGBColor(0x4A, 0x90, 0xD9)
C_DEPT      = RGBColor(0xE0, 0x4F, 0x96)
C_APPT      = RGBColor(0xDC, 0xA8, 0x10)
C_MEDREC    = RGBColor(0x27, 0xAE, 0x60)
C_PAYMENT   = RGBColor(0xE2, 0x69, 0x9F)

F_HEAD = "Poppins"
F_BODY = "Open Sans"
F_MONO = "Roboto Mono"

PAGE_W, PAGE_H = 13.333, 7.5
MARGIN = 0.55
CONTENT_W = PAGE_W - 2 * MARGIN  # 12.233

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------

def slide_new(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = rect(s, 0, 0, PAGE_W, PAGE_H, fill=bg)
    return s


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    _no_shadow(shp)
    shp.text_frame.paragraphs[0].text = ""
    return shp


def set_alpha(shape, pct):
    """pct = opacity percent (e.g. 12 -> 12% visible)."""
    spPr = shape.fill._xPr
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return
    clr = solid.find(qn('a:srgbClr'))
    if clr is None:
        return
    alpha = clr.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))})
    clr.append(alpha)


def cross(s, x, y, size, color, alpha_pct):
    shp = s.shapes.add_shape(MSO_SHAPE.CROSS, Inches(x), Inches(y), Inches(size), Inches(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    set_alpha(shp, alpha_pct)
    shp.line.fill.background()
    _no_shadow(shp)
    return shp


def circle(s, x, y, d, color, alpha_pct=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if alpha_pct is not None:
        set_alpha(shp, alpha_pct)
    shp.line.fill.background()
    _no_shadow(shp)
    return shp


def _apply_run(r, text, font, size, color, bold, italic=False, tracking=None):
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    if tracking:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking * 100)))


def txt(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts:
    {runs:[(text, {font,size,color,bold,italic,tracking})], align, space_after, space_before, line_spacing}
    """
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if p.get('space_after') is not None:
            para.space_after = Pt(p['space_after'])
        else:
            para.space_after = Pt(0)
        if p.get('space_before') is not None:
            para.space_before = Pt(p['space_before'])
        if p.get('line_spacing'):
            para.line_spacing = p['line_spacing']
        for (t, st) in p['runs']:
            r = para.add_run()
            _apply_run(r, t,
                       st.get('font', F_BODY), st.get('size', 12),
                       st.get('color', SLATE), st.get('bold', False),
                       st.get('italic', False), st.get('tracking'))
    return tb


def one(s, x, y, w, h, text, font=F_BODY, size=12, color=SLATE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=None, italic=False,
        line_spacing=None, wrap=True):
    return txt(s, x, y, w, h,
               [{'runs': [(text, {'font': font, 'size': size, 'color': color,
                                  'bold': bold, 'italic': italic, 'tracking': tracking})],
                 'align': align, 'line_spacing': line_spacing}],
               anchor=anchor, wrap=wrap)


IMG_SIZES = {}

def pic(s, name, x, y, w=None, h=None, border=None, border_w=1.0):
    path = os.path.join(A, name)
    if path not in IMG_SIZES:
        IMG_SIZES[path] = Image.open(path).size
    iw, ih = IMG_SIZES[path]
    ar = iw / ih
    if w is not None and h is None:
        h = w / ar
    elif h is not None and w is None:
        w = h * ar
    p = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if border is not None:
        p.line.color.rgb = border
        p.line.width = Pt(border_w)
    else:
        p.line.fill.background()
    _no_shadow(p)
    return p, w, h


def pill(s, x, y, w, h, text, fill, color, size=10.5, bold=True, font=F_BODY,
         tracking=1.0, line=None):
    shp = rect(s, x, y, w, h, fill=fill, line=line, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    _apply_run(r, text, font, size, color, bold, tracking=tracking)
    return shp


def header(s, num, title, presenter, page, total=8, dark=False):
    title_c = WHITE if dark else INK
    num_c = TEAL_BRIGHT if dark else TEAL
    # section number
    one(s, MARGIN, 0.44, 1.0, 0.55, num, font=F_HEAD, size=30, color=num_c, bold=True, wrap=False)
    nw = 0.62 if len(num) <= 2 else 0.9
    one(s, MARGIN + nw + 0.16, 0.50, 8.4, 0.5, title, font=F_HEAD, size=24,
        color=title_c, bold=True, tracking=0.6, wrap=False)
    rect(s, MARGIN + nw + 0.18, 1.08, 0.85, 0.055, fill=AMBER)
    # presenter pill (right)
    label = "PRESENTERS — " + presenter if ("&" in presenter) else "PRESENTER — " + presenter
    pw = max(1.9, 0.36 + 0.092 * len(label))
    pill(s, PAGE_W - MARGIN - pw, 0.52, pw, 0.42, label.upper(),
         fill=(NAVY2 if dark else TEAL_LIGHT), color=(TEAL_BRIGHT if dark else RGBColor(0x0B, 0x66, 0x5C)),
         size=10, tracking=0.9)
    footer(s, page, total, dark=dark)


def footer(s, page, total=8, dark=False):
    c = RGBColor(0x5C, 0x74, 0x90) if dark else SLATE_LIGHT
    rect(s, MARGIN, 7.13, CONTENT_W, 0.012, fill=(NAVY2 if dark else BORDER))
    one(s, MARGIN, 7.2, 6.0, 0.25, "Clinic Appointment Management System",
        size=8.5, color=c, tracking=0.5)
    one(s, PAGE_W - MARGIN - 1.5, 7.2, 1.5, 0.25, f"{page:02d} / {total:02d}",
        size=8.5, color=c, align=PP_ALIGN.RIGHT, tracking=0.5)


KW = {'font': F_MONO, 'size': 11, 'color': TEAL_BRIGHT, 'bold': True}
ID = {'font': F_MONO, 'size': 11, 'color': CODE_ID}
ST = {'font': F_MONO, 'size': 11, 'color': CODE_STR}


def code_card(s, x, y, w, lines, size=11, pad=0.17, line_spacing=1.32):
    n = len(lines)
    h = pad * 2 + n * (size / 72.0) * line_spacing + 0.06
    card = rect(s, x, y, w, h, fill=NAVY, radius=0.10)
    paras = []
    for ln in lines:
        runs = []
        for (t, kind) in ln:
            base = {'kw': KW, 'id': ID, 'str': ST}[kind]
            st = dict(base)
            st['size'] = size
            runs.append((t, st))
        paras.append({'runs': runs, 'line_spacing': line_spacing})
    txt(s, x + pad + 0.05, y + pad, w - 2 * (pad + 0.02), h - 2 * pad, paras)
    return h


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide_new(bg=NAVY)
circle(s, 9.0, -2.4, 5.4, TEAL, 10)
circle(s, -1.6, 5.4, 3.6, TEAL, 8)
circle(s, 7.6, 6.3, 1.7, AMBER, 12)
cross(s, 1.1, 1.0, 0.42, TEAL_BRIGHT, 16)
cross(s, 7.3, 0.8, 0.3, WHITE, 10)
cross(s, 2.1, 6.3, 0.3, WHITE, 9)
cross(s, 12.3, 5.6, 0.42, TEAL_BRIGHT, 14)

pill(s, 0.85, 1.30, 2.9, 0.44, "DBMS GROUP PROJECT", fill=NAVY2, color=TEAL_BRIGHT,
     size=10.5, tracking=1.6)
txt(s, 0.85, 1.98, 7.65, 1.9, [
    {'runs': [("Clinic Appointment", {'font': F_HEAD, 'size': 40, 'color': WHITE, 'bold': True})],
     'line_spacing': 1.06},
    {'runs': [("Management ", {'font': F_HEAD, 'size': 40, 'color': WHITE, 'bold': True}),
              ("System", {'font': F_HEAD, 'size': 40, 'color': TEAL_BRIGHT, 'bold': True})],
     'line_spacing': 1.06},
])
rect(s, 0.88, 3.50, 1.35, 0.075, fill=AMBER)
one(s, 0.88, 3.82, 7.0, 1.3,
    "A relational database that manages patients, doctors, appointments, "
    "medical records and payments — replacing paper records with one "
    "reliable, consistent system.",
    size=13.5, color=MIST, line_spacing=1.45)
one(s, 0.88, 6.55, 7.0, 0.4, "LibreOffice Base (.odb)   ·   Embedded HSQLDB   ·   Normalized to 2NF",
    size=11, color=TEAL_BRIGHT, bold=True, tracking=0.4)

# team & agenda panel
panel = rect(s, 8.62, 1.30, 4.18, 5.62, fill=NAVY2, radius=0.055)
one(s, 8.92, 1.56, 3.6, 0.3, "TEAM  ·  AGENDA", font=F_HEAD, size=11.5,
    color=TEAL_BRIGHT, bold=True, tracking=2.0)
rows = [
    ("01", "Introduction", "Sultan"),
    ("02", "Requirements", "Elbek & Allisa"),
    ("03", "ER Diagram", "Aza"),
    ("04", "Database Demo", "KX"),
    ("05", "SQL Queries", "Abdushukur"),
    ("06", "Conclusion", "Jeremy"),
]
ry = 2.02
for num, sec, who in rows:
    one(s, 8.92, ry + 0.05, 0.5, 0.3, num, font=F_MONO, size=11, color=AMBER, bold=True)
    one(s, 9.42, ry, 2.4, 0.3, sec, font=F_HEAD, size=12.5, color=WHITE, bold=True, wrap=False)
    one(s, 9.42, ry + 0.285, 3.2, 0.26, who, size=10, color=MIST, wrap=False)
    if num != "06":
        rect(s, 8.92, ry + 0.625, 3.58, 0.01, fill=RGBColor(0x1B, 0x47, 0x73))
    ry += 0.79
one(s, 0.88, 7.05, 7.0, 0.3, "Database Management Systems — Group Assignment · June 2026",
    size=9.5, color=RGBColor(0x5C, 0x74, 0x90))

# =====================================================================
# SLIDE 2 — 01 INTRODUCTION (Sultan)
# =====================================================================
s = slide_new()
header(s, "01", "INTRODUCTION", "Sultan", 2)

# problem card
rect(s, MARGIN, 1.46, 5.92, 4.1, fill=WHITE, line=BORDER, radius=0.045)
one(s, 0.85, 1.72, 3.0, 0.3, "THE PROBLEM", font=F_HEAD, size=12, color=CORAL,
    bold=True, tracking=1.6)
one(s, 0.85, 2.06, 5.4, 0.32, "Paper records, forms & spreadsheets", font=F_HEAD,
    size=14.5, color=INK, bold=True, wrap=False)
problems = [
    "Double bookings — no central booking record",
    "Time wasted searching through physical files",
    "Records get lost or damaged",
    "Communication breakdowns when schedules change",
    "Reports become a paper-based nightmare",
]
py = 2.62
for p in problems:
    circle(s, 0.88, py + 0.085, 0.09, CORAL)
    one(s, 1.12, py, 5.1, 0.45, p, size=11.5, color=SLATE)
    py += 0.55

# solution card
rect(s, 6.77, 1.46, 6.0, 4.1, fill=WHITE, line=BORDER, radius=0.045)
one(s, 7.07, 1.72, 3.0, 0.3, "THE SOLUTION", font=F_HEAD, size=12, color=TEAL,
    bold=True, tracking=1.6)
one(s, 7.07, 2.06, 5.5, 0.32, "One centralized relational database", font=F_HEAD,
    size=14.5, color=INK, bold=True, wrap=False)
solutions = [
    ("Patient information", "registration & personal details"),
    ("Doctor information", "specializations & departments"),
    ("Appointments", "date, time, attending doctor, status"),
    ("Medical records", "diagnosis, treatment, medicine"),
    ("Payments", "amount, method, status, date"),
]
sy = 2.62
for t, d in solutions:
    circle(s, 7.10, sy + 0.085, 0.09, TEAL)
    txt(s, 7.34, sy, 5.25, 0.45, [{'runs': [
        (t, {'size': 11.5, 'color': INK, 'bold': True}),
        ("  —  " + d, {'size': 11, 'color': SLATE}),
    ]}])
    sy += 0.55

# benefits band
rect(s, MARGIN, 5.78, CONTENT_W, 1.12, fill=NAVY, radius=0.07)
benefits = ["Improves efficiency", "Reduces data duplication", "Ensures data consistency"]
bx = 1.25
bw = 3.9
for i, b in enumerate(benefits):
    txt(s, bx + i * bw, 6.02, bw, 0.4, [{'runs': [
        ("✓  ", {'font': F_BODY, 'size': 14, 'color': TEAL_BRIGHT, 'bold': True}),
        (b, {'font': F_HEAD, 'size': 13.5, 'color': WHITE, 'bold': True}),
    ], 'align': PP_ALIGN.CENTER}])
one(s, MARGIN, 6.46, CONTENT_W, 0.3,
    "Built with LibreOffice Base (.odb) — staff find what they need in seconds, not minutes",
    size = 10, color=MIST, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 3 — 02 REQUIREMENTS (Elbek & Allisa)
# =====================================================================
s = slide_new()
header(s, "02", "REQUIREMENTS — SIX MAIN ENTITIES", "Elbek & Allisa", 3)

cards = [
    ("1", C_PATIENT, "Patient", [("PK", "PatientID")],
     "Personal details of every patient — name, gender, DOB, phone, address.", "11 rows"),
    ("2", C_DOCTOR, "Doctor", [("PK", "DoctorID"), ("FK", "DepartmentID")],
     "Doctor details & specialization, each linked to a department.", "6 rows"),
    ("3", C_DEPT, "Department", [("PK", "DepartmentID")],
     "Clinic departments with their name and location in the building.", "6 rows"),
    ("4", C_APPT, "Appointment", [("PK", "AppointmentID"), ("FK", "PatientID · DoctorID")],
     "Links patients & doctors — date, time, reason and status of each visit.", "10 rows"),
    ("5", C_MEDREC, "MedicalRecord", [("PK", "RecordID"), ("FK", "AppointmentID")],
     "Diagnosis, treatment and medicine recorded after each completed visit.", "5 rows"),
    ("6", C_PAYMENT, "Payment", [("PK", "PaymentID"), ("FK", "AppointmentID")],
     "Billing per appointment — amount, method, status and payment date.", "5 rows"),
]
cw, ch, gap = 3.95, 2.24, 0.19
for i, (num, color, name, keys, desc, rowsn) in enumerate(cards):
    cx = MARGIN + (i % 3) * (cw + gap)
    cy = 1.52 + (i // 3) * (ch + 0.2)
    rect(s, cx, cy, cw, ch, fill=WHITE, line=BORDER, radius=0.06)
    chip = rect(s, cx + 0.22, cy + 0.20, 0.36, 0.36, fill=color, radius=0.28)
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    pr = tf.paragraphs[0]
    pr.alignment = PP_ALIGN.CENTER
    rr = pr.add_run()
    _apply_run(rr, num, F_HEAD, 13, WHITE, True)
    one(s, cx + 0.72, cy + 0.225, 2.6, 0.34, name, font=F_HEAD, size=14.5,
        color=INK, bold=True, wrap=False)
    ky = cy + 0.68
    for klabel, kfield in keys:
        txt(s, cx + 0.24, ky, cw - 0.45, 0.26, [{'runs': [
            (klabel + "  ", {'font': F_MONO, 'size': 9, 'color': (AMBER if klabel == "PK" else TEAL), 'bold': True}),
            (kfield, {'font': F_MONO, 'size': 9, 'color': SLATE}),
        ]}])
        ky += 0.245
    one(s, cx + 0.24, cy + 1.22, cw - 0.48, 0.62, desc, size=9.8, color=SLATE,
        line_spacing=1.22)
    pill(s, cx + cw - 1.12, cy + ch - 0.42, 0.9, 0.27, rowsn, fill=CHIP_GRAY,
         color=INK, size=8.5, tracking=0.4)

# users strip
rect(s, MARGIN, 6.42, CONTENT_W, 0.58, fill=NAVY, radius=0.10)
txt(s, MARGIN + 0.35, 6.56, 8.6, 0.32, [{'runs': [
    ("USER ROLES   ", {'font': F_HEAD, 'size': 10, 'color': TEAL_BRIGHT, 'bold': True, 'tracking': 1.4}),
    ("Receptionist / Admin  ·  Doctor  ·  Clinic Manager  ·  Patient",
     {'size': 11.5, 'color': WHITE, 'bold': True}),
]}])
one(s, 7.6, 6.56, 5.0, 0.32, "Every table stores at least 5 sample records",
    size=10, color=MIST, align=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 4 — 03 ER DIAGRAM (Aza)
# =====================================================================
s = slide_new()
header(s, "03", "ER DIAGRAM & RELATIONSHIPS", "Aza", 4)

rels = [
    ("Department", "Doctor", "1 : N"),
    ("Patient", "Appointment", "1 : N"),
    ("Doctor", "Appointment", "1 : N"),
    ("Appointment", "MedicalRecord", "1 : 1"),
    ("Appointment", "Payment", "1 : N"),
]
ry = 1.58
for a, b, card_ in rels:
    rect(s, MARGIN, ry, 4.3, 0.56, fill=WHITE, line=BORDER, radius=0.16)
    pill(s, MARGIN + 0.14, ry + 0.115, 0.78, 0.33, card_, fill=TEAL_LIGHT,
         color=RGBColor(0x0B, 0x66, 0x5C), size=10, font=F_MONO, tracking=0.2)
    txt(s, MARGIN + 1.06, ry + 0.135, 3.2, 0.3, [{'runs': [
        (a, {'size': 11.5, 'color': INK, 'bold': True}),
        ("  →  ", {'size': 11.5, 'color': TEAL, 'bold': True}),
        (b, {'size': 11.5, 'color': INK, 'bold': True}),
    ]}])
    ry += 0.71

note = rect(s, MARGIN, 5.22, 4.3, 1.62, fill=TEAL_LIGHT, radius=0.07)
one(s, MARGIN + 0.22, 5.40, 3.9, 0.3, "APPOINTMENT IS THE CENTRAL ENTITY",
    font=F_HEAD, size=10.5, color=RGBColor(0x0B, 0x66, 0x5C), bold=True, tracking=1.0)
one(s, MARGIN + 0.22, 5.74, 3.92, 1.0,
    "Patients make appointments and doctors handle them — medical records "
    "and payments always link back to an appointment.",
    size=10.5, color=RGBColor(0x11, 0x4E, 0x46), line_spacing=1.3)

# ERD image card
rect(s, 5.05, 1.46, 7.73, 5.38, fill=WHITE, line=BORDER, radius=0.04)
ew = 7.33
p, w, h = pic(s, "erd-lucidchart.png", 5.25, 1.66, w=ew)
one(s, 5.05, 1.66 + h + 0.12, 7.73, 0.3,
    "ER diagram in crow's foot notation — created with Lucidchart",
    size=9.5, color=SLATE_LIGHT, align=PP_ALIGN.CENTER, italic=True)

# =====================================================================
# SLIDE 5 — 04 DATABASE DEMO (KX)
# =====================================================================
s = slide_new()
header(s, "04", "DATABASE DEMO — LIBREOFFICE BASE", "KX", 5)

# relationships screenshot card
rect(s, MARGIN, 1.46, 6.45, 3.45, fill=WHITE, line=BORDER, radius=0.045)
one(s, 0.80, 1.62, 5.9, 0.3, "Relationships view — foreign keys in action",
    font=F_HEAD, size=12, color=INK, bold=True)
pic(s, "base-relationships.png", 0.80, 1.98, w=5.95, border=BORDER)

# implementation facts card
rect(s, 7.18, 1.46, 5.6, 3.45, fill=WHITE, line=BORDER, radius=0.045)
one(s, 7.46, 1.62, 5.0, 0.3, "Implementation", font=F_HEAD, size=12, color=INK, bold=True)
facts = [
    ("LibreOffice Base 24.2", "with the embedded HSQLDB 1.8 engine"),
    (".odb file format", "no separate database server required"),
    ("2NF design", "all six tables built from the normalized model"),
    ("Referential integrity", "enforced through foreign keys"),
    ("Realistic sample data", "every table holds at least 5 records"),
]
fy = 2.04
for t, d in facts:
    circle(s, 7.49, fy + 0.082, 0.09, TEAL)
    txt(s, 7.72, fy, 4.95, 0.5, [{'runs': [
        (t, {'size': 11, 'color': INK, 'bold': True}),
        ("  —  " + d, {'size': 10.5, 'color': SLATE}),
    ], 'line_spacing': 1.1}])
    fy += 0.565

# three table screenshots
shots = [
    ("table-patient.png", "Patient — 11 rows"),
    ("table-doctor.png", "Doctor — 6 rows"),
    ("table-appointment.png", "Appointment — 10 rows"),
]
sh = 1.40
widths = []
for f, _ in shots:
    iw, ih = Image.open(os.path.join(A, f)).size
    widths.append(sh * iw / ih)
gap2 = 0.18
total = sum(widths) + 2 * gap2
sx = MARGIN + (CONTENT_W - total) / 2
for (f, label), w_ in zip(shots, widths):
    one(s, sx, 5.12, w_, 0.26, label, font=F_HEAD, size=10, color=INK, bold=True,
        align=PP_ALIGN.CENTER, wrap=False)
    pic(s, f, sx, 5.40, w=w_, border=BORDER)
    sx += w_ + gap2

# =====================================================================
# SLIDE 6 — 05 SQL QUERIES (1/2) (Abdushukur)
# =====================================================================
s = slide_new()
header(s, "05", "SQL QUERIES — SELECT & JOIN", "Abdushukur", 6)

# Query 1
pill(s, MARGIN, 1.50, 0.62, 0.34, "Q1", fill=TEAL, color=WHITE, size=11, tracking=0.6)
one(s, 1.32, 1.51, 5.2, 0.32, "View All Patients", font=F_HEAD, size=15, color=INK,
    bold=True, wrap=False)
one(s, MARGIN, 1.95, 5.5, 0.55, "Retrieves all 11 patient records — the full Patient table.",
    size=11, color=SLATE)
code_card(s, MARGIN, 2.42, 5.5, [
    [("SELECT", "kw"), (" * ", "id"), ("FROM", "kw"), (' "Patient"', "id")],
], size=12)
one(s, MARGIN, 3.30, 5.5, 0.3, "11 of 11 records returned", size=9.5,
    color=SLATE_LIGHT, italic=True)
pic(s, "query1-view-all-patients.png", 7.62, 1.50, w=4.62, border=BORDER)

rect(s, MARGIN, 4.12, CONTENT_W, 0.012, fill=BORDER)

# Query 2
pill(s, MARGIN, 4.30, 0.62, 0.34, "Q2", fill=TEAL, color=WHITE, size=11, tracking=0.6)
one(s, 1.32, 4.31, 5.2, 0.32, "Appointment Details (JOIN)", font=F_HEAD, size=15,
    color=INK, bold=True, wrap=False)
one(s, MARGIN, 4.74, 5.45, 0.55,
    "Joins Appointment with Patient and Doctor — names instead of IDs.",
    size=11, color=SLATE)
code_card(s, MARGIN, 5.20, 5.5, [
    [("SELECT", "kw"), (' "A"."AppointmentID", "P"."PatientName",', "id")],
    [('       "D"."DoctorName", "A"."AppointmentDate",', "id")],
    [('       "A"."AppointmentTime", "A"."Status"', "id")],
    [("FROM", "kw"), (' "Appointment" "A"', "id")],
    [("JOIN", "kw"), (' "Patient" "P" ', "id"), ("ON", "kw"), (' "A"."PatientID" = "P"."PatientID"', "id")],
    [("JOIN", "kw"), (' "Doctor" "D" ', "id"), ("ON", "kw"), (' "A"."DoctorID" = "D"."DoctorID"', "id")],
], size=9.5)
pic(s, "query2-appointment-details.png", 6.35, 4.78, w=6.43, border=BORDER)

# =====================================================================
# SLIDE 7 — 05 SQL QUERIES (2/2) (Abdushukur)
# =====================================================================
s = slide_new()
header(s, "05", "SQL QUERIES — INSERT & UPDATE", "Abdushukur", 7)

# Query 3
pill(s, MARGIN, 1.50, 0.62, 0.34, "Q3", fill=AMBER, color=WHITE, size=11, tracking=0.6)
one(s, 1.32, 1.51, 6.0, 0.32, "Insert New Patient", font=F_HEAD, size=15, color=INK,
    bold=True, wrap=False)
one(s, MARGIN, 1.95, 6.6, 0.32,
    "Adds patient #11 — Alicia Wong — to the Patient table.", size=11, color=SLATE)
h3 = code_card(s, MARGIN, 2.40, 6.75, [
    [("INSERT INTO", "kw"), (' "Patient" ("PatientID", "PatientName",', "id")],
    [('  "Gender", "DOB", "PhoneNumber", "Address")', "id")],
    [("VALUES", "kw"), (" (11, ", "id"), ("'Alicia Wong'", "str"), (", ", "id"),
     ("'Female'", "str"), (", ", "id"), ("'2002-08-15'", "str"), (",", "id")],
    [("  ", "id"), ("'0135558888'", "str"), (", ", "id"),
     ("'Austin Heights, Johor Bahru, Johor'", "str"), (");", "id")],
], size=10)
one(s, 7.62, 1.62, 5.0, 0.26, "Query Design view — LibreOffice Base", size=9.5,
    color=SLATE_LIGHT, italic=True, align=PP_ALIGN.CENTER)
pic(s, "query3-insert-new-patient.png", 7.62, 1.92, w=5.16, border=BORDER)
one(s, 7.62, 3.1, 5.0, 0.6, "After running: the Patient table grows from 10 to 11 rows.",
    size=10, color=SLATE, align=PP_ALIGN.CENTER)

rect(s, MARGIN, 4.12, CONTENT_W, 0.012, fill=BORDER)

# Query 4
pill(s, MARGIN, 4.30, 0.62, 0.34, "Q4", fill=AMBER, color=WHITE, size=11, tracking=0.6)
one(s, 1.32, 4.31, 6.0, 0.32, "Update Appointment Status", font=F_HEAD, size=15,
    color=INK, bold=True, wrap=False)
one(s, MARGIN, 4.74, 6.6, 0.32,
    "Changes the status of appointment No. 5 via its AppointmentID.",
    size=11, color=SLATE)
code_card(s, MARGIN, 5.20, 6.75, [
    [("UPDATE", "kw"), (' "Appointment"', "id")],
    [("SET", "kw"), (' "Status" = ', "id"), ("'Scheduled'", "str")],
    [("WHERE", "kw"), (' "AppointmentID" = 5;', "id")],
], size=11)
one(s, 7.62, 4.55, 5.0, 0.26, "Query Design view — LibreOffice Base", size=9.5,
    color=SLATE_LIGHT, italic=True, align=PP_ALIGN.CENTER)
pic(s, "query4-update-status.png", 7.85, 4.88, w=4.7, border=BORDER)
one(s, 7.62, 6.0, 5.0, 0.6, "UPDATE + WHERE — only the matching row changes.",
    size=10, color=SLATE, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 8 — 06 CONCLUSION (Jeremy)
# =====================================================================
s = slide_new(bg=NAVY)
circle(s, 10.6, -2.0, 4.6, TEAL, 9)
circle(s, -1.4, 5.6, 3.2, TEAL, 8)
cross(s, 12.45, 5.9, 0.36, TEAL_BRIGHT, 14)
cross(s, 0.7, 1.0, 0.3, WHITE, 9)
header(s, "06", "CONCLUSION", "Jeremy", 8, dark=True)

txt(s, MARGIN, 1.75, 12.2, 1.5, [
    {'runs': [("The Clinic Appointment Management System successfully manages ",
               {'font': F_HEAD, 'size': 19, 'color': WHITE, 'bold': True}),
              ("patients, doctors, appointments, medical records and payments.",
               {'font': F_HEAD, 'size': 19, 'color': TEAL_BRIGHT, 'bold': True})],
     'line_spacing': 1.3},
])
one(s, MARGIN, 2.95, 12.0, 0.5,
    "The database is normalized to 2NF, keeps data consistent through foreign keys, "
    "and supports SQL operations efficiently.",
    size=13, color=MIST, line_spacing=1.35)

stats = [("6", "TABLES"), ("6", "DEPARTMENTS"), ("6", "DOCTORS"),
         ("11", "PATIENTS"), ("10", "APPOINTMENTS")]
sw, sgap = 2.32, 0.155
sx = MARGIN
for n, label in stats:
    rect(s, sx, 3.85, sw, 1.32, fill=NAVY2, radius=0.09)
    one(s, sx, 4.05, sw, 0.55, n, font=F_HEAD, size=28, color=TEAL_BRIGHT, bold=True,
        align=PP_ALIGN.CENTER)
    one(s, sx, 4.72, sw, 0.3, label, size=9.5, color=MIST, bold=True,
        align=PP_ALIGN.CENTER, tracking=1.2)
    sx += sw + sgap

one(s, MARGIN, 5.55, 12.2, 0.9, "Thank you!", font=F_HEAD, size=36, color=WHITE,
    bold=True, align=PP_ALIGN.CENTER)
one(s, MARGIN, 6.42, 12.2, 0.4, "Questions are welcome.", size=12.5, color=MIST,
    align=PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT, os.path.getsize(OUT) // 1024, "KB")
